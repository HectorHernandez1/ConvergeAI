import hashlib
import warnings
from pathlib import Path
from typing import Optional, Tuple
import PyPDF2
import pdfplumber
import pymupdf as fitz
from pptx import Presentation
from config import settings

def extract_text(file_path: str) -> str:
    """Extract text from PDF or PPT/PPTX file."""
    path = Path(file_path)
    suffix = path.suffix.lower()

    # Route to appropriate extractor based on file type
    if suffix == '.pdf':
        text = _try_pdf_extraction(file_path)
        if text and len(text.strip()) >= 100:
            return text
    elif suffix in ['.ppt', '.pptx']:
        text = _try_ppt_extraction(file_path)
        if text and len(text.strip()) >= 100:
            return text
    else:
        # Unknown file type - try both methods as fallback
        text = _try_pdf_extraction(file_path)
        if text and len(text.strip()) >= 100:
            return text
        text = _try_ppt_extraction(file_path)
        if text and len(text.strip()) >= 100:
            return text

    raise ValueError(f"Failed to extract meaningful text from {file_path}. "
                     "Ensure the file is not a scanned image-only document or corrupted.")

def _try_pdf_extraction(pdf_path: str) -> Optional[str]:
    """Try to extract text from PDF using fallback chain."""
    text = _try_pypdf2(pdf_path)
    
    if text and len(text.strip()) >= 100:
        return text
    
    text = _try_pdfplumber(pdf_path)
    
    if text and len(text.strip()) >= 100:
        return text
    
    text = _try_pymupdf(pdf_path)
    
    if text and len(text.strip()) >= 100:
        return text
    
    return None

def _try_ppt_extraction(ppt_path: str) -> Optional[str]:
    """Try to extract text from PPT/PPTX file."""
    path = Path(ppt_path)
    suffix = path.suffix.lower()

    if suffix not in ['.ppt', '.pptx']:
        return None

    try:
        text_parts = []

        if suffix == '.pptx':
            prs = Presentation(ppt_path)

            for slide in prs.slides:
                for shape in slide.shapes:
                    # Extract text from regular text shapes
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text.strip())

                    # Extract text from tables (accessed via shape.has_table)
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text:
                                    row_text.append(cell.text.strip())
                            if row_text:
                                text_parts.append(" | ".join(row_text))

            return "\n".join(filter(None, text_parts)) if text_parts else None

        elif suffix == '.ppt':
            warnings.warn(f"PPT format (not PPTX) has limited text extraction support: {ppt_path}")
            return f"PPT file: {path.name} (Note: PPT format has limited text extraction - consider converting to PPTX)"

        return None
    except ImportError:
        return None
    except Exception:
        return None

def _try_pypdf2(pdf_path: str) -> Optional[str]:
    try:
        text_parts = []
        with open(pdf_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text_parts.append(extracted)
        return "\n".join(text_parts) if text_parts else None
    except Exception:
        return None

def _try_pdfplumber(pdf_path: str) -> Optional[str]:
    try:
        text_parts = []
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                extracted = page.extract_text()
                if extracted:
                    text_parts.append(extracted)
        return "\n".join(text_parts) if text_parts else None
    except Exception:
        return None

def _try_pymupdf(pdf_path: str) -> Optional[str]:
    try:
        doc = fitz.open(pdf_path)
        text_parts = []
        for page in doc:
            text_parts.append(page.get_text())
        doc.close()
        return "\n".join(text_parts) if text_parts else None
    except Exception:
        return None

def get_pdf_hash(pdf_path: str) -> str:
    with open(pdf_path, 'rb') as f:
        return hashlib.md5(f.read()).hexdigest()
