import hashlib
import mimetypes
import warnings
from pathlib import Path
from typing import Optional, Tuple
import PyPDF2
import pdfplumber
import pymupdf as fitz
from pptx import Presentation
from bs4 import BeautifulSoup
from config import settings
from utils.image_types import ExtractedContent, ExtractedImage

IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.gif', '.webp', '.bmp'}

def extract_text(file_path: str) -> str:
    """Extract text from PDF, PPT/PPTX, or HTML file."""
    path = Path(file_path)
    suffix = path.suffix.lower()

    # Route to appropriate extractor based on file type
    if suffix == '.pdf':
        text = _try_pdf_extraction(file_path)
        if text and len(text.strip()) >= 100:
            return text
    elif suffix in ['.ppt', '.pptx']:
        text = _try_ppt_extraction(file_path)
        if text and len(text.strip()) >= 100:
            return text
    elif suffix in ['.html', '.htm']:
        text = _try_html_extraction(file_path)
        if text and len(text.strip()) >= 100:
            return text
    else:
        # Unknown file type - try all methods as fallback
        text = _try_pdf_extraction(file_path)
        if text and len(text.strip()) >= 100:
            return text
        text = _try_ppt_extraction(file_path)
        if text and len(text.strip()) >= 100:
            return text
        text = _try_html_extraction(file_path)
        if text and len(text.strip()) >= 100:
            return text

    raise ValueError(f"Failed to extract meaningful text from {file_path}. "
                     "Ensure the file is not a scanned image-only document or corrupted.")


def extract_content(file_path: str) -> ExtractedContent:
    """Extract text and images from a supported file.

    For PDFs: extracts text (existing chain) + embedded images via pymupdf.
    For image files (PNG/JPG/etc.): returns empty text with the image.
    For PPT/HTML: extracts text only (no image extraction).
    """
    path = Path(file_path)
    suffix = path.suffix.lower()

    # Direct image file
    if suffix in IMAGE_EXTENSIONS:
        return _load_image_file(file_path)

    # PDF: extract text + images
    if suffix == '.pdf':
        text = _try_pdf_extraction(file_path) or ""
        images = _extract_pdf_images(file_path)
        if not text.strip() and not images:
            raise ValueError(f"Failed to extract content from {file_path}. "
                             "Ensure the file is not a scanned image-only document or corrupted.")
        return ExtractedContent(text=text, images=images)

    # PPT/HTML: text only for now
    if suffix in ['.ppt', '.pptx']:
        text = _try_ppt_extraction(file_path) or ""
    elif suffix in ['.html', '.htm']:
        text = _try_html_extraction(file_path) or ""
    else:
        text = ""

    if not text.strip():
        raise ValueError(f"Failed to extract meaningful text from {file_path}. "
                         "Ensure the file is not a scanned image-only document or corrupted.")
    return ExtractedContent(text=text, images=[])


def _extract_pdf_images(pdf_path: str) -> list[ExtractedImage]:
    """Extract embedded images from PDF using pymupdf.

    Skips tiny images (< 5KB) which are typically icons/decorations.
    """
    max_images = settings.max_images
    images = []
    try:
        doc = fitz.open(pdf_path)
        for page_num, page in enumerate(doc):
            if len(images) >= max_images:
                break
            image_list = page.get_images(full=True)
            for img_index, img_info in enumerate(image_list):
                if len(images) >= max_images:
                    break
                xref = img_info[0]
                base_image = doc.extract_image(xref)
                if base_image and len(base_image["image"]) >= 5000:
                    ext = base_image["ext"]
                    media_type = f"image/{ext}" if ext != "jpg" else "image/jpeg"
                    images.append(ExtractedImage(
                        data=base_image["image"],
                        media_type=media_type,
                        source_page=page_num + 1,
                        label=f"Page {page_num + 1}, Image {img_index + 1}"
                    ))
        doc.close()
    except Exception:
        pass
    return images


def _load_image_file(file_path: str) -> ExtractedContent:
    """Load a standalone image file as ExtractedContent with no text."""
    path = Path(file_path)
    mime_type, _ = mimetypes.guess_type(file_path)
    if not mime_type or not mime_type.startswith("image/"):
        mime_type = "image/png"

    with open(file_path, "rb") as f:
        data = f.read()

    image = ExtractedImage(
        data=data,
        media_type=mime_type,
        label=path.name
    )
    return ExtractedContent(text="", images=[image])


def _try_pdf_extraction(pdf_path: str) -> Optional[str]:
    """Try to extract text from PDF using fallback chain."""
    text = _try_pypdf2(pdf_path)
    
    if text and len(text.strip()) >= 100:
        return text
    
    text = _try_pdfplumber(pdf_path)
    
    if text and len(text.strip()) >= 100:
        return text
    
    text = _try_pymupdf(pdf_path)
    
    if text and len(text.strip()) >= 100:
        return text
    
    return None

def _try_ppt_extraction(ppt_path: str) -> Optional[str]:
    """Try to extract text from PPT/PPTX file."""
    path = Path(ppt_path)
    suffix = path.suffix.lower()

    if suffix not in ['.ppt', '.pptx']:
        return None

    try:
        text_parts = []

        if suffix == '.pptx':
            prs = Presentation(ppt_path)

            for slide in prs.slides:
                for shape in slide.shapes:
                    # Extract text from regular text shapes
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text.strip())

                    # Extract text from tables (accessed via shape.has_table)
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text:
                                    row_text.append(cell.text.strip())
                            if row_text:
                                text_parts.append(" | ".join(row_text))

            return "\n".join(filter(None, text_parts)) if text_parts else None

        elif suffix == '.ppt':
            warnings.warn(f"PPT format (not PPTX) has limited text extraction support: {ppt_path}")
            return f"PPT file: {path.name} (Note: PPT format has limited text extraction - consider converting to PPTX)"

        return None
    except ImportError:
        return None
    except Exception:
        return None

def _try_html_extraction(html_path: str) -> Optional[str]:
    """Try to extract text from HTML file."""
    try:
        with open(html_path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        
        # Get text and clean it up
        text = soup.get_text()
        
        # Clean up whitespace
        lines = [line.strip() for line in text.splitlines()]
        text_parts = [line for line in lines if line]
        
        return "\n".join(text_parts) if text_parts else None
    except Exception:
        return None

def _try_pypdf2(pdf_path: str) -> Optional[str]:
    try:
        text_parts = []
        with open(pdf_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text_parts.append(extracted)
        return "\n".join(text_parts) if text_parts else None
    except Exception:
        return None

def _try_pdfplumber(pdf_path: str) -> Optional[str]:
    try:
        text_parts = []
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                extracted = page.extract_text()
                if extracted:
                    text_parts.append(extracted)
        return "\n".join(text_parts) if text_parts else None
    except Exception:
        return None

def _try_pymupdf(pdf_path: str) -> Optional[str]:
    try:
        doc = fitz.open(pdf_path)
        text_parts = []
        for page in doc:
            text_parts.append(page.get_text())
        doc.close()
        return "\n".join(text_parts) if text_parts else None
    except Exception:
        return None

def get_pdf_hash(pdf_path: str) -> str:
    with open(pdf_path, 'rb') as f:
        return hashlib.md5(f.read()).hexdigest()
